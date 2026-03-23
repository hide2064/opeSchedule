"""
opeSchedule 運用マニュアル PowerPoint 生成スクリプト
実行: python docs/make_manual_pptx.py
出力: docs/opeSchedule_manual.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── カラーパレット ────────────────────────────────────────────────────────────
C_PRIMARY   = RGBColor(0x4A, 0x90, 0xD9)   # #4A90D9 ブランドブルー
C_DARK      = RGBColor(0x2C, 0x3E, 0x50)   # ダークネイビー
C_LIGHT_BG  = RGBColor(0xF5, 0xF6, 0xFA)   # 薄グレー背景
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_ACCENT    = RGBColor(0x27, 0xAE, 0x60)   # グリーン（成功）
C_WARN      = RGBColor(0xE7, 0x4C, 0x3C)   # レッド（注意）
C_MUTED     = RGBColor(0x7F, 0x8C, 0x8D)   # グレー（補足）
C_SECTION   = RGBColor(0x2F, 0x72, 0xB8)   # セクションヘッダー

# スライドサイズ: 16:9 ワイド
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # 完全ブランクレイアウト


# ── ヘルパー関数 ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=14, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, space_before=6, bullet=False, indent=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """共通ヘッダーバー（上部 0.8inch）"""
    add_rect(slide, 0, 0, 13.33, 0.8, C_PRIMARY)
    add_text(slide, "📅  opeSchedule  運用マニュアル",
             0.2, 0.05, 8, 0.65, size=16, bold=True, color=C_WHITE)
    add_text(slide, title,
             0.2, 0.85, 12.8, 0.6, size=26, bold=True, color=C_PRIMARY)
    if subtitle:
        add_text(slide, subtitle,
                 0.2, 1.5, 12.8, 0.45, size=15, color=C_MUTED)
    # 下部アクセントライン
    add_rect(slide, 0, 7.3, 13.33, 0.2, C_PRIMARY)


def icon_box(slide, icon, label, x, y, w=2.4, h=1.6,
             bg=C_PRIMARY, fg=C_WHITE):
    """アイコン + ラベル のカードブロック"""
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, icon, x, y + 0.1, w, 0.7,
             size=28, align=PP_ALIGN.CENTER, color=fg)
    add_text(slide, label, x, y + 0.85, w, 0.6,
             size=13, bold=True, align=PP_ALIGN.CENTER, color=fg)


def step_box(slide, num, text, x, y, w=3.8, h=0.75,
             bg=C_PRIMARY, fg=C_WHITE):
    """ステップ番号付きボックス"""
    add_rect(slide, x, y, 0.55, h, bg)
    add_text(slide, str(num), x + 0.01, y + 0.08, 0.55, 0.55,
             size=20, bold=True, align=PP_ALIGN.CENTER, color=fg)
    add_rect(slide, x + 0.55, y, w - 0.55, h, C_LIGHT_BG)
    add_text(slide, text, x + 0.65, y + 0.12, w - 0.75, 0.55,
             size=13, color=C_DARK)


def info_card(slide, title, lines, x, y, w=4.0, h=2.0,
              title_bg=C_PRIMARY, body_bg=C_LIGHT_BG):
    """タイトル付き情報カード"""
    add_rect(slide, x, y, w, 0.4, title_bg)
    add_text(slide, title, x + 0.1, y + 0.04, w - 0.2, 0.35,
             size=13, bold=True, color=C_WHITE)
    add_rect(slide, x, y + 0.4, w, h - 0.4, body_bg)
    txb = slide.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.45),
        Inches(w - 0.24), Inches(h - 0.55))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = C_DARK


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 1: タイトル
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# 背景グラデーション風（2段）
add_rect(sl, 0, 0, 13.33, 5.0, C_PRIMARY)
add_rect(sl, 0, 5.0, 13.33, 2.5, C_DARK)

# タイトル
add_text(sl, "📅  opeSchedule", 1.0, 1.0, 11.3, 1.2,
         size=54, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "運 用 マ ニ ュ ア ル", 1.0, 2.3, 11.3, 0.9,
         size=32, bold=False, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)
add_text(sl, "Web ベース開発スケジュール管理ツール（ガントチャート）",
         1.0, 3.2, 11.3, 0.55, size=16, color=RGBColor(0xCC, 0xDD, 0xFF),
         align=PP_ALIGN.CENTER)

# バージョン・日付
add_rect(sl, 4.5, 4.1, 4.33, 0.65, RGBColor(0x1A, 0x5C, 0x99))
add_text(sl, "Version 0.1.0  ／  2026-03-20",
         4.5, 4.15, 4.33, 0.55, size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

# 機能アイコン
icons = [("📋", "スケジュール\n管理"), ("📊", "ガントチャート"),
         ("🔄", "バージョン\n管理"), ("📥", "Import /\nExport")]
for i, (ico, lbl) in enumerate(icons):
    icon_box(sl, ico, lbl, 1.5 + i * 2.65, 5.4, 2.3, 1.6,
             bg=RGBColor(0x3A, 0x5C, 0x7A), fg=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 2: 目次
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "目次")

sections = [
    ("1", "システム概要",       "機能概要・アクセス方法"),
    ("2", "Top 画面の操作",     "プロジェクト一覧・作成・設定"),
    ("3", "Schedule 画面の操作","ガントチャート・タスク管理"),
    ("4", "バージョン管理",     "履歴・バージョンUP・過去版閲覧"),
    ("5", "Import / Export",    "データの入出力"),
    ("6", "比較・フィルター表示","複数プロジェクト閲覧"),
    ("7", "管理者向け情報",     "起動方法・環境変数・マイグレーション"),
    ("8", "よくあるトラブル",   "エラー対応・ブラウザキャッシュ"),
]

cols = [(0.4, 3.6), (6.8, 6.1)]
rows_per_col = 4
for i, (num, title, desc) in enumerate(sections):
    col = i // rows_per_col
    row = i %  rows_per_col
    cx, cy_base = cols[col][0], 2.05 + row * 1.2
    cw = cols[col][1]

    add_rect(sl, cx, cy_base, 0.7, 0.85, C_PRIMARY)
    add_text(sl, num, cx + 0.01, cy_base + 0.1, 0.7, 0.65,
             size=22, bold=True, align=PP_ALIGN.CENTER, color=C_WHITE)
    add_rect(sl, cx + 0.7, cy_base, cw - 0.7, 0.85, C_WHITE)
    add_text(sl, title, cx + 0.85, cy_base + 0.04, cw - 1.0, 0.42,
             size=16, bold=True, color=C_PRIMARY)
    add_text(sl, desc,  cx + 0.85, cy_base + 0.46, cw - 1.0, 0.35,
             size=12, color=C_MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 3: システム概要
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "1. システム概要", "opeSchedule とは何か")

# 左: 概要テキスト
add_rect(sl, 0.3, 1.9, 6.0, 5.0, C_WHITE)
add_text(sl, "概要", 0.4, 1.95, 5.8, 0.4, size=14, bold=True, color=C_PRIMARY)

txb = sl.shapes.add_textbox(Inches(0.45), Inches(2.45), Inches(5.6), Inches(4.2))
txb.word_wrap = True
tf = txb.text_frame; tf.word_wrap = True
items = [
    "● Web ブラウザで動作する開発スケジュール管理ツール",
    "● ガントチャート形式でプロジェクトを可視化",
    "● 大項目・中項目・小項目の3階層でタスクを整理",
    "● マイルストーン（◆）で重要イベントを管理",
    "● 進捗率をバーで視覚的に表現",
    "● タスクのドラッグ＆ドロップで日程変更",
    "● 複数プロジェクトの比較・横断表示",
    "● バージョン管理で過去のスケジュールを閲覧",
]
first = True
for item in items:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run(); run.text = item
    run.font.size = Pt(13); run.font.color.rgb = C_DARK

# 右: アクセス情報
info_card(sl, "アクセス URL", [
    "http://localhost:8000",
    "",
    "（本番環境は管理者に確認）",
], 6.8, 1.9, 6.0, 1.6)

info_card(sl, "動作環境", [
    "● 推奨ブラウザ: Google Chrome / Edge",
    "● 画面解像度: 1280px 以上推奨",
    "● ネットワーク: サーバーと同一 LAN",
], 6.8, 3.65, 6.0, 1.8)

info_card(sl, "技術スタック", [
    "Frontend : React 18 + Vite 5",
    "Backend  : Python FastAPI",
    "DB       : SQLite（開発）/ PostgreSQL（本番）",
], 6.8, 5.6, 6.0, 1.6)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 4: Top 画面の構成
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "2. Top 画面の操作①", "画面構成・プロジェクト一覧")

# 画面レイアウト模式図
add_rect(sl, 0.3, 1.9, 12.7, 5.3, C_WHITE)

# ヘッダー模示
add_rect(sl, 0.35, 1.95, 12.6, 0.55, C_PRIMARY)
add_text(sl, "📅  opeSchedule  ヘッダー",
         0.45, 2.0, 12.4, 0.45, size=13, bold=True, color=C_WHITE)

# サイドバー模示
add_rect(sl, 0.35, 2.5, 2.5, 4.65, RGBColor(0xEC, 0xF0, 0xF1))
add_text(sl, "左サイドバー",    0.45, 2.55, 2.3, 0.35, size=11, bold=True, color=C_DARK)
add_text(sl, "▶ 比較表示\n  （折りたたみ）\n\n▶ 大項目フィルター\n  （折りたたみ）",
         0.45, 2.95, 2.3, 1.8, size=11, color=C_DARK)

# プロジェクト一覧
add_rect(sl, 2.9, 2.5, 6.5, 4.65, RGBColor(0xF8, 0xF9, 0xFA))
add_text(sl, "Projects タブ", 3.0, 2.55, 6.3, 0.35, size=11, bold=True, color=C_DARK)

# ヘッダー行
add_rect(sl, 2.9, 2.95, 6.5, 0.38, RGBColor(0xD6, 0xE4, 0xF5))
for col_text, cx in [("色", 2.95), ("プロジェクト名", 3.2), ("状態", 5.2),
                      ("Ver", 6.4), ("最終更新", 7.1), ("操作", 8.6)]:
    add_text(sl, col_text, cx, 2.98, 1.2, 0.32, size=10, bold=True, color=C_PRIMARY)

# サンプル行
rows = [
    ("■", "ECサイトリニューアル", "作業中", "v3", "03/20", "✎ 🗑"),
    ("■", "クラウド移行プロジェクト", "未開始", "v1", "03/18", "✎ 🗑"),
    ("■", "セキュリティ監査", "終了", "—", "02/28", "✎ 🗑"),
]
for ri, (dot, name, status, ver, date, ops) in enumerate(rows):
    ry = 3.38 + ri * 0.5
    bg = C_WHITE if ri % 2 == 0 else RGBColor(0xF8, 0xF9, 0xFA)
    add_rect(sl, 2.9, ry, 6.5, 0.48, bg)
    add_text(sl, dot,   2.93, ry+0.1, 0.25, 0.3, size=14, color=C_PRIMARY)
    add_text(sl, name,  3.2,  ry+0.1, 2.0,  0.3, size=11, color=C_PRIMARY)
    add_text(sl, status,5.2,  ry+0.1, 1.1,  0.3, size=10, color=C_DARK)
    add_text(sl, ver,   6.4,  ry+0.1, 0.65, 0.3, size=11, bold=True, color=C_ACCENT)
    add_text(sl, date,  7.1,  ry+0.1, 1.2,  0.3, size=10, color=C_MUTED)
    add_text(sl, ops,   8.6,  ry+0.1, 0.75, 0.3, size=11, color=C_WARN)

# グローバル設定タブ
add_rect(sl, 9.45, 2.5, 3.6, 4.65, RGBColor(0xF8, 0xF9, 0xFA))
add_text(sl, "Global Config タブ", 9.55, 2.55, 3.4, 0.35, size=11, bold=True, color=C_DARK)
add_text(sl, "・週の開始曜日\n・デフォルト表示モード\n・テーマ（ライト/ダーク）\n・今日へ自動スクロール\n・週末ハイライト",
         9.55, 2.95, 3.4, 2.0, size=11, color=C_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 5: プロジェクト作成・編集
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "2. Top 画面の操作②", "プロジェクトの作成・編集・アーカイブ")

# 作成手順
add_rect(sl, 0.3, 1.9, 5.8, 5.3, C_WHITE)
add_text(sl, "プロジェクト作成手順", 0.4, 1.95, 5.6, 0.4,
         size=14, bold=True, color=C_PRIMARY)
steps = [
    "「+ New Project」ボタンをクリック",
    "プロジェクト名を入力（必須）",
    "表示色・客先名・ベースPJ を任意入力",
    "業務ステータスを選択（未開始 / 作業中 / 中断 / 終了）",
    "「作成」ボタンをクリックして保存",
    "プロジェクト名をクリックして Schedule 画面へ移動",
]
for i, s in enumerate(steps):
    step_box(sl, i + 1, s, 0.4, 2.45 + i * 0.77, w=5.6)

# 右: フィールド説明
info_card(sl, "プロジェクト設定項目", [
    "プロジェクト名 ★ 必須",
    "説明           任意テキスト",
    "表示色         HEX カラー（#4A90D9 等）",
    "業務ステータス 未開始 / 作業中 / 中断 / 終了",
    "客先名         顧客企業名",
    "ベースPJ       参照元プロジェクト名",
    "表示モード     Day/Week/Month/Quarter",
], 6.5, 1.9, 6.5, 3.1)

# アーカイブルール
add_rect(sl, 6.5, 5.15, 6.5, 2.0, RGBColor(0xFF, 0xF3, 0xCD))
add_rect(sl, 6.5, 5.15, 6.5, 0.4, RGBColor(0xF3, 0x9C, 0x12))
add_text(sl, "⚡  アーカイブ自動化ルール",
         6.6, 5.18, 6.3, 0.35, size=13, bold=True, color=C_WHITE)
add_text(sl, "業務ステータスを「中断」または「終了」に変更すると\n"
             "自動的にアーカイブ状態（非表示）になります。\n"
             "Top 画面の「アーカイブ済みを表示 □」で再表示可能。",
         6.6, 5.6, 6.3, 1.45, size=12, color=C_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 6: Schedule 画面の構成
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "3. Schedule 画面の操作①", "ガントチャートの画面構成")

# ヘッダーバー模示
add_rect(sl, 0.3, 1.9, 12.7, 0.6, C_PRIMARY)
btns = [("← Top", 0.4), ("📅 opeSchedule", 1.4), ("PJ名", 3.4),
        ("[Day][Week][Month][Quarter]", 5.5), ("+ Add  JSON  CSV  📋履歴[2]", 9.8)]
for txt, bx in btns:
    add_text(sl, txt, bx, 2.0, 2.8, 0.42, size=10, color=C_WHITE)

# 階層ペイン
add_rect(sl, 0.3, 2.55, 4.2, 4.6, RGBColor(0xEC, 0xF0, 0xF1))
add_text(sl, "左ペイン（HierarchyPane）",
         0.35, 2.58, 4.1, 0.35, size=11, bold=True, color=C_DARK)
for ci, col in enumerate(["大項目", "中項目", "小項目（タスク）"]):
    cx = 0.35 + ci * 1.4
    add_rect(sl, cx, 2.98, 1.38, 0.32, C_PRIMARY)
    add_text(sl, col, cx + 0.03, 3.01, 1.33, 0.28,
             size=10, bold=True, color=C_WHITE)

for ri in range(5):
    ry = 3.35 + ri * 0.62
    add_rect(sl, 0.35, ry, 1.38, 0.58, RGBColor(0xD6, 0xE4, 0xF5) if ri < 3 else C_LIGHT_BG)
    add_rect(sl, 1.75, ry, 1.38, 0.58, RGBColor(0xE8, 0xF4, 0xFD) if ri in (0,3) else C_LIGHT_BG)
    add_rect(sl, 3.13, ry, 1.35, 0.58, C_WHITE)

add_text(sl, "Phase1\n要件定義", 0.4, 3.37, 1.28, 1.18, size=10, color=C_DARK)
add_text(sl, "調査・\n要件整理", 1.8, 3.37, 1.28, 1.18, size=10, color=C_DARK)
add_text(sl, "市場調査\nヒアリング\n要件書作成\n設計\n◆ 完了",
         3.18, 3.37, 1.25, 3.0, size=9, color=C_DARK)

# ガントペイン
add_rect(sl, 4.55, 2.55, 8.45, 4.6, C_WHITE)
add_text(sl, "右ペイン（GanttPane）",
         4.6, 2.58, 8.3, 0.35, size=11, bold=True, color=C_DARK)
add_rect(sl, 4.55, 2.98, 8.45, 0.32, RGBColor(0xD6, 0xE4, 0xF5))
for di, d in enumerate(["4/1", "4/8", "4/15", "4/22", "5/1", "5/8", "5/15"]):
    add_text(sl, d, 4.6 + di * 1.18, 3.01, 1.15, 0.28, size=9, color=C_DARK)

# タスクバー
bars = [(0, 0, 3, C_PRIMARY), (1, 2, 2, C_ACCENT), (2, 1, 4, C_PRIMARY),
        (3, 4, 2, C_ACCENT), (4, 3, 0, C_WARN)]
for ri, sx, sw, bc in bars:
    ry = 3.35 + ri * 0.62
    if sw == 0:  # マイルストーン
        add_text(sl, "◆", 4.6 + sx * 1.18 - 0.1, ry + 0.15, 0.4, 0.35,
                 size=16, color=bc, align=PP_ALIGN.CENTER)
    else:
        add_rect(sl, 4.6 + sx * 1.18, ry + 0.12, sw * 1.18, 0.35, bc)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 7: タスク操作
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "3. Schedule 画面の操作②", "タスクの追加・編集・日程変更")

# タスク追加
info_card(sl, "① タスク追加", [
    "1. 「+ Add Task」ボタンをクリック",
    "2. 大項目・中項目・小項目名を入力",
    "3. 開始日・終了日を設定",
    "4. 種別選択: task（期間）/ milestone（◆）",
    "5. 「作成」ボタンで保存",
], 0.3, 1.9, 6.2, 2.5)

# タスク編集
info_card(sl, "② タスク詳細・編集", [
    "1. タスクバーまたは左ペインの小項目をクリック",
    "2. 右隣に詳細パネルがポップアップ",
    "3. 名前・日程・進捗率・色・メモを編集",
    "4. 依存関係タスクを設定",
    "5. 「保存」ボタンで確定",
], 6.7, 1.9, 6.3, 2.5)

# 日程変更
info_card(sl, "③ ドラッグ＆ドロップで日程変更", [
    "1. タスクバーにマウスを合わせる",
    "2. 左右にドラッグすると日程がシフト",
    "3. マウスを離した時点で自動保存",
    "※ マイルストーン（◆）はドラッグ不可",
    "※ 比較モード・履歴モードでは操作不可",
], 0.3, 4.6, 6.2, 2.55)

# 表示モード
info_card(sl, "④ 表示モード切替", [
    "ヘッダーの [Day][Week][Month][Quarter] で切替",
    "",
    "Day    : 1px = 40px   詳細表示",
    "Week   : 1px =  8px   週単位",
    "Month  : 1px =  2.5px 月単位",
    "Quarter: 1px =  0.8px 四半期",
], 6.7, 4.6, 6.3, 2.55)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 8: バージョン管理
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "4. バージョン管理", "履歴の記録・バージョンUP・過去版の閲覧")

# フロー図
add_rect(sl, 0.3, 1.9, 12.7, 1.3, C_WHITE)
add_text(sl, "変更ログ → バージョンUP の流れ",
         0.4, 1.95, 12.5, 0.4, size=13, bold=True, color=C_PRIMARY)

flow_items = [
    ("タスク\n変更操作", C_PRIMARY),
    ("変更ログ\n自動記録", C_ACCENT),
    ("📋 履歴\nボタン [n]", C_DARK),
    ("バージョン\nUP ボタン", C_WARN),
    ("スナップ\nショット保存", RGBColor(0x8E, 0x44, 0xAD)),
]
for i, (label, color) in enumerate(flow_items):
    bx = 0.5 + i * 2.55
    add_rect(sl, bx, 2.4, 2.2, 0.7, color)
    add_text(sl, label, bx, 2.45, 2.2, 0.65,
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "→", bx + 2.2, 2.62, 0.35, 0.35,
                 size=18, bold=True, color=C_DARK)

# 3カラム説明
for col_x, col_title, col_color, col_items in [
    (0.3, "変更ログ（自動）", C_ACCENT, [
        "タスクを操作するたびに DB へ",
        "自動記録される軽量ログ",
        "",
        "記録される操作:",
        "  ＋ タスク追加",
        "  ✎ タスク更新",
        "  ✕ タスク削除",
        "  📅 日程変更",
        "  ↕ 並び替え",
    ]),
    (4.7, "バージョンUP（手動）", C_PRIMARY, [
        "① 「📋 履歴」ボタンをクリック",
        "② 未コミット変更一覧を確認",
        "③ バージョン名を入力（任意）",
        "   例: v3 リリース準備完了",
        "④ 「⬆ バージョンUP」をクリック",
        "",
        "→ 現在のタスク全量をスナップ",
        "  ショットとして保存",
        "→ 最大 50 バージョン保持",
    ]),
    (9.1, "過去版の閲覧", RGBColor(0x8E, 0x44, 0xAD), [
        "① 履歴パネルの過去バージョン",
        "   一覧をクリック",
        "② ガントチャートが過去の",
        "   状態に切り替わる",
        "③ 読み取り専用（編集不可）",
        "④ ヘッダーに版数と名称表示",
        "   例: v2 — 初回納品",
        "⑤「現在に戻る」で最新に戻る",
    ]),
]:
    add_rect(sl, col_x, 3.3, 4.0, 0.4, col_color)
    add_text(sl, col_title, col_x + 0.1, 3.33, 3.8, 0.35,
             size=13, bold=True, color=C_WHITE)
    add_rect(sl, col_x, 3.7, 4.0, 3.5, C_WHITE)
    txb = sl.shapes.add_textbox(
        Inches(col_x + 0.12), Inches(3.78), Inches(3.8), Inches(3.3))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for line in col_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run(); run.text = line
        run.font.size = Pt(12); run.font.color.rgb = C_DARK


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 9: Import / Export
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "5. Import / Export", "データの入出力")

# Export
add_rect(sl, 0.3, 1.9, 6.2, 5.3, C_WHITE)
add_rect(sl, 0.3, 1.9, 6.2, 0.45, C_ACCENT)
add_text(sl, "📤  Export（エクスポート）", 0.4, 1.93, 6.0, 0.4,
         size=14, bold=True, color=C_WHITE)

txb = sl.shapes.add_textbox(Inches(0.4), Inches(2.42), Inches(6.0), Inches(4.6))
txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
export_lines = [
    ("手順:", True),
    ("① Schedule 画面ヘッダーの「JSON」または「CSV」ボタンをクリック", False),
    ("② ファイルが自動ダウンロードされます", False),
    ("", False),
    ("JSON 形式:", True),
    ("・プロジェクト情報 + タスク全件 + 依存関係", False),
    ("・ファイル名: project_{id}.json", False),
    ("・再インポート可能なフルバックアップ", False),
    ("", False),
    ("CSV 形式:", True),
    ("・タスクデータのみ（プロジェクト情報なし）", False),
    ("・ファイル名: project_{id}.csv", False),
    ("・Excel 等での加工に便利", False),
    ("", False),
    ("※ 比較モード・履歴モードでは Export 不可", False),
]
first = True
for text, bold in export_lines:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(5 if bold else 3)
    run = p.add_run(); run.text = text
    run.font.size = Pt(12); run.font.bold = bold
    run.font.color.rgb = C_PRIMARY if bold else C_DARK

# Import
add_rect(sl, 6.8, 1.9, 6.2, 5.3, C_WHITE)
add_rect(sl, 6.8, 1.9, 6.2, 0.45, C_PRIMARY)
add_text(sl, "📥  Import（インポート）", 6.9, 1.93, 6.0, 0.4,
         size=14, bold=True, color=C_WHITE)

txb2 = sl.shapes.add_textbox(Inches(6.9), Inches(2.42), Inches(6.0), Inches(4.6))
txb2.word_wrap = True; tf2 = txb2.text_frame; tf2.word_wrap = True
import_lines = [
    ("手順:", True),
    ("① Schedule 画面ヘッダーの「Import」ボタンをクリック", False),
    ("② JSON または CSV ファイルを選択", False),
    ("③ 新規プロジェクトとして登録されます", False),
    ("", False),
    ("対応形式:", True),
    ("JSON: opeSchedule 形式（フルバックアップ）", False),
    ("CSV:  タスクのみ。プロジェクト名はファイル名から取得", False),
    ("", False),
    ("制限事項:", True),
    ("・ファイルサイズ上限: 10 MB", False),
    ("・循環依存（A→B→A）は自動検出でエラー", False),
    ("・既存プロジェクトへの追加ではなく", False),
    ("  新規プロジェクトとして作成される", False),
    ("", False),
    ("サンプルデータ: docs/sample_*.json を利用可", False),
]
first = True
for text, bold in import_lines:
    p = tf2.paragraphs[0] if first else tf2.add_paragraph()
    first = False
    p.space_before = Pt(5 if bold else 3)
    run = p.add_run(); run.text = text
    run.font.size = Pt(12); run.font.bold = bold
    run.font.color.rgb = C_PRIMARY if bold else C_DARK


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 10: 比較・フィルター表示
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "6. 比較・フィルター表示", "複数プロジェクトの横断閲覧")

# 比較表示
add_rect(sl, 0.3, 1.9, 6.2, 5.3, C_WHITE)
add_rect(sl, 0.3, 1.9, 6.2, 0.45, C_PRIMARY)
add_text(sl, "📊  比較表示（Compare View）", 0.4, 1.93, 6.0, 0.4,
         size=14, bold=True, color=C_WHITE)

info_card(sl, "操作手順", [
    "① Top 画面左サイドバー「比較表示」を展開",
    "② 比較したいプロジェクトを 2つ以上チェック",
    "③「まとめて表示」ボタンをクリック",
    "④ ガントチャートに全 PJ のタスクが並んで表示",
], 0.4, 2.45, 6.0, 2.1, title_bg=C_MUTED)

info_card(sl, "URL 形式", [
    "/schedule?projects=1,2,3",
    "（カンマ区切りで project ID を指定）",
], 0.4, 4.65, 6.0, 1.0, title_bg=C_MUTED)

add_rect(sl, 0.4, 5.75, 6.0, 1.3, RGBColor(0xFF, 0xF3, 0xCD))
add_text(sl, "⚠  比較モードの制限",
         0.5, 5.78, 5.9, 0.35, size=12, bold=True, color=C_WARN)
add_text(sl, "・タスク追加・Export ボタン非表示\n"
             "・タスク詳細は読み取り専用（編集不可）\n"
             "・ドラッグ・クリティカルパス・依存矢印無効",
         0.5, 6.16, 5.9, 0.85, size=11, color=C_DARK)

# フィルター表示
add_rect(sl, 6.8, 1.9, 6.2, 5.3, C_WHITE)
add_rect(sl, 6.8, 1.9, 6.2, 0.45, RGBColor(0x8E, 0x44, 0xAD))
add_text(sl, "🔍  大項目フィルター表示", 6.9, 1.93, 6.0, 0.4,
         size=14, bold=True, color=C_WHITE)

info_card(sl, "操作手順", [
    "① Top 画面左サイドバー「大項目フィルター」を展開",
    "② 表示したい大項目をチェック（複数可）",
    "③ 対象プロジェクトをチェック",
    "④「フィルター表示」ボタンをクリック",
    "⑤ 選択した大項目のタスクのみが表示される",
], 6.9, 2.45, 6.0, 2.3, title_bg=RGBColor(0x8E, 0x44, 0xAD))

info_card(sl, "URL 形式", [
    "/schedule?projects=1&projects=2",
    "           &catfilter=Phase1要件定義",
], 6.9, 4.85, 6.0, 1.0, title_bg=C_MUTED)

add_rect(sl, 6.9, 5.95, 6.0, 1.1, RGBColor(0xEA, 0xF4, 0xFB))
add_text(sl, "💡  活用例",
         7.0, 5.98, 5.9, 0.35, size=12, bold=True, color=C_PRIMARY)
add_text(sl, "「Phase1 要件定義」フェーズのみ全プロジェクト横断で\n"
             "進捗を比較する際に便利です。",
         7.0, 6.35, 5.9, 0.65, size=11, color=C_DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 11: 管理者向け情報
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "7. 管理者向け情報", "起動・環境変数・DB マイグレーション")

# 起動手順（Windows）
info_card(sl, "ローカル起動（Windows）", [
    "start.bat をダブルクリック or コマンドから実行",
    "",
    "[1/4] pip install -r requirements.txt",
    "[2/4] alembic upgrade head",
    "[3/4] npm install && npm run build",
    "[4/4] uvicorn app.main:app --reload ポート 8000",
    "",
    "→ http://localhost:8000 でアクセス",
], 0.3, 1.9, 6.2, 3.2)

# 環境変数
info_card(sl, "主要環境変数（backend/.env）", [
    "DATABASE_URL  = sqlite:///./opeschedule.db",
    "APP_ENV       = development",
    "CORS_ORIGINS  = http://localhost:8000",
    "LOG_LEVEL     = info",
    "APP_WORKERS   = 1",
    "",
    "本番: DATABASE_URL を PostgreSQL に変更",
], 6.8, 1.9, 6.2, 3.2)

# Docker
info_card(sl, "Docker 起動（PostgreSQL 込み）", [
    "docker-compose up",
    "",
    "→ PostgreSQL + FastAPI が起動",
    "→ alembic upgrade head が自動実行",
    "→ http://localhost:8000 でアクセス",
], 0.3, 5.25, 6.2, 1.9)

# マイグレーション
info_card(sl, "DB マイグレーション", [
    "# 最新に適用",
    "alembic upgrade head",
    "",
    "# 新規マイグレーション作成",
    "alembic revision --autogenerate -m '説明'",
    "",
    "CI: push/PR で ruff + pytest + docker build",
], 6.8, 5.25, 6.2, 1.9)


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 12: よくあるトラブル
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "8. よくあるトラブル", "エラー対応・ブラウザキャッシュ")

troubles = [
    ("画面が古い表示のまま",
     "ブラウザキャッシュが残っています。",
     "Ctrl + Shift + R（強制リロード）を実施。\n"
     "それでも改善しない場合は DevTools の\n"
     "Network タブで読み込まれているファイルのハッシュ値を確認。"),
    ("「履歴の読み込みエラー: Not Found」",
     "バックエンドのルーター登録漏れまたはサーバーが古いコードで動作中。",
     "サーバーを再起動してください。\n"
     "Windows: すべての Python プロセスを終了後、\n"
     "start.bat を再実行。"),
    ("「バージョンUP エラー: Method Not Allowed」",
     "POST /snapshots エンドポイントが未登録のサーバーに接続中。",
     "上記と同様にサーバーを再起動してください。"),
    ("タスクバーのドラッグが効かない",
     "比較モードまたは履歴モードでは編集操作が無効になっています。",
     "URL を確認し、単一プロジェクトモード\n"
     "（?project=<id>）でアクセスしてください。\n"
     "履歴モードの場合は「現在に戻る」をクリック。"),
    ("インポートに失敗する",
     "循環依存（A→B→A）またはファイルサイズ超過（10MB）の可能性。",
     "依存関係を見直してください。\n"
     "エラーメッセージに詳細が表示されます。"),
    ("プロジェクトが一覧に表示されない",
     "アーカイブ状態（status=archived）になっています。",
     "「アーカイブ済みを表示 □」チェックボックスを\n"
     "オンにして確認してください。"),
]

for i, (symptom, cause, solution) in enumerate(troubles):
    col = i % 2
    row = i // 2
    bx = 0.3 + col * 6.5
    by = 1.9 + row * 1.85

    add_rect(sl, bx, by, 6.2, 1.75, C_WHITE)
    add_rect(sl, bx, by, 6.2, 0.38, C_WARN)
    add_text(sl, f"⚠  {symptom}", bx + 0.1, by + 0.03, 6.0, 0.34,
             size=12, bold=True, color=C_WHITE)

    txb = sl.shapes.add_textbox(
        Inches(bx + 0.1), Inches(by + 0.43), Inches(6.0), Inches(1.25))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = f"原因: {cause}"
    run.font.size = Pt(10); run.font.color.rgb = C_MUTED; run.font.italic = True

    p2 = tf.add_paragraph(); p2.space_before = Pt(4)
    run2 = p2.add_run(); run2.text = f"対処: {solution}"
    run2.font.size = Pt(11); run2.font.color.rgb = C_DARK


# ═══════════════════════════════════════════════════════════════════════════════
# スライド 13: まとめ / 操作早見表
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_LIGHT_BG)
slide_header(sl, "操作早見表", "よく使う操作のショートカット")

ops = [
    ("プロジェクト作成",   "Top 画面 → 「+ New Project」"),
    ("スケジュール表示",   "Top 画面 → プロジェクト名をクリック"),
    ("タスク追加",         "Schedule 画面 → 「+ Add Task」"),
    ("タスク編集",         "タスクバー / 左ペインの小項目をクリック"),
    ("日程変更（D&D）",   "タスクバーを左右にドラッグ"),
    ("表示モード切替",     "ヘッダーの [Day][Week][Month][Quarter]"),
    ("履歴パネル開く",     "ヘッダーの「📋 履歴」ボタン"),
    ("バージョンUP",       "履歴パネル → バージョン名入力 → 「⬆ バージョンUP」"),
    ("過去版表示",         "履歴パネル → 過去バージョン一覧をクリック"),
    ("現在に戻る",         "ヘッダーの履歴モードバナー → 「現在に戻る」"),
    ("JSON エクスポート",  "ヘッダーの「JSON」ボタン"),
    ("CSV エクスポート",   "ヘッダーの「CSV」ボタン"),
    ("インポート",         "ヘッダーの「Import」ボタン"),
    ("比較表示",           "Top 画面サイドバー → 比較表示 → チェック → まとめて表示"),
    ("フィルター表示",     "Top 画面サイドバー → 大項目フィルター → チェック → フィルター表示"),
    ("アーカイブ",         "プロジェクト Edit → project_status を「中断」または「終了」に変更"),
]

cols = 2
rows_per_col = len(ops) // cols + (1 if len(ops) % cols else 0)
for i, (action, how) in enumerate(ops):
    col = i // rows_per_col
    row = i %  rows_per_col
    bx = 0.3 + col * 6.5
    by = 1.95 + row * 0.33

    add_rect(sl, bx, by, 2.6, 0.3, C_PRIMARY if row % 2 == 0 else RGBColor(0x2F, 0x72, 0xB8))
    add_text(sl, action, bx + 0.08, by + 0.02, 2.5, 0.27,
             size=11, bold=True, color=C_WHITE)
    add_rect(sl, bx + 2.6, by, 3.6, 0.3, C_WHITE if row % 2 == 0 else C_LIGHT_BG)
    add_text(sl, how, bx + 2.7, by + 0.02, 3.5, 0.27,
             size=10.5, color=C_DARK)


# ── 保存 ──────────────────────────────────────────────────────────────────────
out = "docs/opeSchedule_manual.pptx"
prs.save(out)
print(f"[OK] Saved: {out}")
print(f"     Slides: {len(prs.slides)}")
