"""
opeSchedule 管理職向け提案資料 PowerPoint 生成スクリプト
実行: python docs/make_proposal_pptx.py
出力: docs/opeSchedule_proposal.pptx

内容:
  1. 表紙
  2. 目次
  3. 現状の課題 (PowerApps)
  4. 提案システム概要 (opeSchedule)
  5. PowerApps vs opeSchedule 機能比較
  6. 運用コスト比較 (PowerApps / AWS / オンプレ)
  7. AWS 構成案
  8. オンプレ構成案
  9. コスト試算サマリー
 10. 移行・導入スケジュール
 11. リスクと対策
 12. 推奨構成と意思決定フロー
 13. まとめ・次のアクション
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラーパレット ────────────────────────────────────────────────────────────
C_PRIMARY  = RGBColor(0x1A, 0x56, 0x96)  # ディープブルー（信頼感）
C_ACCENT   = RGBColor(0x0D, 0x9E, 0x6A)  # グリーン（メリット・OK）
C_WARN     = RGBColor(0xD6, 0x3B, 0x2F)  # レッド（課題・コスト）
C_GOLD     = RGBColor(0xE8, 0x8C, 0x0A)  # ゴールド（推奨・注目）
C_DARK     = RGBColor(0x1C, 0x2B, 0x3A)  # ダークネイビー（テキスト）
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)  # 薄ブルー背景
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED    = RGBColor(0x6B, 0x7C, 0x8D)  # グレー（補足）
C_POWERAPPS= RGBColor(0x74, 0x27, 0xC8)  # PowerApps パープル
C_AWS      = RGBColor(0xFF, 0x99, 0x00)  # AWS オレンジ
C_ONPREM   = RGBColor(0x2E, 0x86, 0xAB)  # オンプレ ブルー

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── ヘルパー関数 ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             size=14, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=13, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, space_before=5):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, label="opeSchedule 導入提案"):
    """共通ヘッダーバー（上部 0.8inch）"""
    add_rect(slide, 0, 0, 13.33, 0.75, C_PRIMARY)
    add_text(slide, label, 0.2, 0.08, 9, 0.6,
             size=15, bold=True, color=C_WHITE)
    add_text(slide, "管理職向け", 11.5, 0.08, 1.7, 0.6,
             size=13, bold=True, color=C_GOLD, align=PP_ALIGN.RIGHT)
    add_text(slide, title, 0.2, 0.85, 12.8, 0.65,
             size=26, bold=True, color=C_PRIMARY)
    if subtitle:
        add_text(slide, subtitle, 0.2, 1.52, 12.8, 0.42,
                 size=14, color=C_MUTED)
    add_rect(slide, 0, 7.3, 13.33, 0.2, C_PRIMARY)


def cost_bar(slide, label, val_str, ratio, x, y, w_max=7.0,
             bar_color=C_PRIMARY, h=0.55):
    """横棒コスト比較バー"""
    add_text(slide, label, x, y + 0.05, 2.3, 0.45,
             size=13, bold=True, color=C_DARK)
    bar_w = max(w_max * ratio, 0.2)
    add_rect(slide, x + 2.4, y, bar_w, h, bar_color)
    add_text(slide, val_str, x + 2.4 + bar_w + 0.1, y + 0.08,
             2.5, 0.42, size=13, bold=True, color=C_DARK)


def compare_row(slide, cols, x_list, y, heights=0.55,
                bg=None, text_color=C_DARK, bold=False, size=12):
    """比較表の1行"""
    if bg:
        add_rect(slide, x_list[0], y,
                 x_list[-1] - x_list[0] + 2.5, heights, bg)
    for i, (col, x) in enumerate(zip(cols, x_list)):
        col_w = (x_list[i+1] - x) if i + 1 < len(x_list) else 2.5
        add_text(slide, col, x + 0.08, y + 0.06,
                 col_w - 0.16, heights - 0.1,
                 size=size, bold=bold, color=text_color, wrap=True)


# ── スライド 1: 表紙 ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)

# 背景グラデーション風（2段矩形）
add_rect(sl, 0, 0, 13.33, 5.0, C_PRIMARY)
add_rect(sl, 0, 5.0, 13.33, 2.5, C_DARK)

# タイトル
add_text(sl, "opeSchedule", 1.0, 0.6, 11.33, 1.1,
         size=54, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "開発スケジュール管理ツール  導入提案資料",
         1.0, 1.75, 11.33, 0.8,
         size=24, bold=False, color=RGBColor(0xB0, 0xCE, 0xF0),
         align=PP_ALIGN.CENTER)

# サブタイトルライン
add_rect(sl, 3.5, 2.65, 6.33, 0.06, C_GOLD)

add_text(sl, "管理職向け  |  コスト・比較・導入計画",
         1.0, 2.8, 11.33, 0.6,
         size=18, color=C_GOLD, align=PP_ALIGN.CENTER)

# 特徴 3 アイコン
for i, (icon, txt) in enumerate([
    ("PowerApps との比較", "現行システムとの\n機能・コスト差を整理"),
    ("AWS / オンプレ 運用案", "2つの構成案と\nコスト試算を提示"),
    ("移行計画", "段階的導入と\nリスク対策を明示"),
]):
    bx = 1.0 + i * 3.9
    add_rect(sl, bx, 3.6, 3.4, 1.5, RGBColor(0x0D, 0x2A, 0x45))
    add_text(sl, icon, bx + 0.1, 3.65, 3.2, 0.55,
             size=13, bold=True, color=C_GOLD)
    add_text(sl, txt, bx + 0.1, 4.22, 3.2, 0.8,
             size=11, color=RGBColor(0xB0, 0xCE, 0xF0))

add_text(sl, "2026年3月  情報システム部門",
         1.0, 5.2, 11.33, 0.5,
         size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
add_text(sl, "本資料は管理職の意思決定を支援するために作成されました",
         1.0, 5.7, 11.33, 0.4,
         size=11, italic=True,
         color=RGBColor(0x70, 0x80, 0x90), align=PP_ALIGN.CENTER)


# ── スライド 2: 目次 ──────────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "目次", "本資料の構成")

toc = [
    ("1", "現状の課題 — PowerApps 運用の問題点"),
    ("2", "opeSchedule 提案システム概要"),
    ("3", "PowerApps vs opeSchedule 機能比較"),
    ("4", "運用コスト比較（PowerApps / AWS / オンプレ）"),
    ("5", "AWS 構成案"),
    ("6", "オンプレ構成案"),
    ("7", "コスト試算サマリー"),
    ("8", "移行・導入スケジュール"),
    ("9", "リスクと対策"),
    ("10", "推奨構成と意思決定フロー"),
    ("11", "まとめ・次のアクション"),
]

for i, (num, text) in enumerate(toc):
    row = i % 6
    col = i // 6
    x = 0.5 + col * 6.6
    y = 2.1 + row * 0.82
    bg = C_LIGHT_BG if i % 2 == 0 else C_WHITE
    add_rect(sl, x, y, 6.2, 0.7, bg)
    add_rect(sl, x, y, 0.55, 0.7, C_PRIMARY)
    add_text(sl, num, x + 0.01, y + 0.1, 0.55, 0.5,
             size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, text, x + 0.65, y + 0.12, 5.4, 0.5,
             size=13, color=C_DARK)

add_rect(sl, 0.5, 7.05, 12.33, 0.06, C_GOLD)


# ── スライド 3: 現状の課題（PowerApps）──────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "現状の課題", "PowerApps で運用している開発スケジュール管理の問題点")

issues = [
    ("ライセンスコスト大", C_WARN,
     "Power Apps per-user プランは 1 人あたり約\n"
     "3,000 円/月。20 名規模で月 60,000 円（年 72 万円）。\n"
     "ユーザー増加に比例してコストが膨らむ。"),
    ("カスタマイズの制約", C_WARN,
     "ガントチャートの細かな表示制御（マイルストーン◆、\n"
     "バー色指定、ドラッグ）が標準では困難。\n"
     "追加開発にはプレミアムコネクタ費用が発生。"),
    ("データ主権・可搬性", C_WARN,
     "データは Microsoft クラウドに保存。\n"
     "エクスポート・移行が困難で、サービス停止リスク。\n"
     "オフライン利用・ローカル保存は不可。"),
    ("バージョン管理の欠如", C_WARN,
     "スケジュールの過去版を参照・比較する\n"
     "機能が標準では提供されない。\n"
     "変更履歴も手動管理が必要。"),
    ("ベンダーロックイン", C_WARN,
     "Microsoft 365 ライセンス体系の変更や\n"
     "Price 改定に全面依存。価格交渉力ゼロ。\n"
     "M365 解約時はシステム即廃止。"),
    ("学習・保守コスト", C_WARN,
     "Power Fx（独自言語）の習得が必要。\n"
     "社内エンジニアでの改修が難しく\n"
     "外部ベンダーへの依存度が高い。"),
]

for i, (title, color, desc) in enumerate(issues):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 2.0 + row * 2.45
    add_rect(sl, x, y, 4.0, 0.45, color)
    add_text(sl, title, x + 0.1, y + 0.05, 3.8, 0.38,
             size=13, bold=True, color=C_WHITE)
    add_rect(sl, x, y + 0.45, 4.0, 1.85, C_LIGHT_BG)
    txb = sl.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.55),
        Inches(3.76), Inches(1.65))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].add_run().text = desc
    tf.paragraphs[0].runs[0].font.size = Pt(11)
    tf.paragraphs[0].runs[0].font.color.rgb = C_DARK

add_text(sl, "現行 PowerApps の主要課題",
         0.3, 1.6, 4.0, 0.35,
         size=12, bold=True, color=C_MUTED, italic=True)


# ── スライド 4: 提案システム概要 ─────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "opeSchedule 提案システム概要",
             "自社開発の Web ベースガントチャート — 拡張・保守・コストすべてに優れる")

# 左: 機能一覧
add_rect(sl, 0.3, 2.0, 6.2, 0.45, C_PRIMARY)
add_text(sl, "主要機能", 0.4, 2.05, 6.0, 0.38,
         size=14, bold=True, color=C_WHITE)

features = [
    "インタラクティブ ガントチャート（ドラッグ & ドロップ）",
    "マイルストーン / 通常タスク の混在管理",
    "手動バージョンUP + 変更ログ（差分追跡）",
    "プロジェクト毎の JSON / CSV インポート・エクスポート",
    "タスク依存関係（矢印線）の可視化",
    "モデル名 による最上位分類（プロジェクト群を整理）",
    "サムネイル画像 をヘッダーに表示（視認性向上）",
    "複数プロジェクト比較表示（セパレーター行で視認性向上）",
    "REST API 完備（他システムとの連携容易）",
]
txb = sl.shapes.add_textbox(
    Inches(0.4), Inches(2.5), Inches(5.9), Inches(4.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
for j, feat in enumerate(features):
    p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
    p.space_before = Pt(4)
    run = p.add_run()
    run.text = "  ●  " + feat
    run.font.size = Pt(12)
    run.font.color.rgb = C_DARK

# 右: 技術スタック
add_rect(sl, 6.8, 2.0, 6.2, 0.45, C_DARK)
add_text(sl, "技術スタック", 6.9, 2.05, 6.0, 0.38,
         size=14, bold=True, color=C_WHITE)

stacks = [
    ("バックエンド", "Python / FastAPI", C_ACCENT),
    ("フロントエンド", "React 18 + Vite", C_PRIMARY),
    ("ガントライブラリ", "Frappe Gantt (OSS)", C_PRIMARY),
    ("データベース", "SQLite (開発) / PostgreSQL (本番)", C_ACCENT),
    ("コンテナ", "Docker / docker-compose", RGBColor(0x09,0x9C,0xEC)),
    ("CI/CD", "GitHub Actions", C_DARK),
]
for i, (label, val, col) in enumerate(stacks):
    y = 2.55 + i * 0.66
    add_rect(sl, 6.8, y, 2.0, 0.55, C_LIGHT_BG)
    add_text(sl, label, 6.9, y + 0.08, 1.9, 0.42, size=11, color=C_MUTED)
    add_rect(sl, 8.8, y, 4.2, 0.55, col)
    add_text(sl, val, 8.9, y + 0.1, 3.9, 0.38,
             size=12, bold=True, color=C_WHITE)

# 下部: ポイント
add_rect(sl, 0.3, 6.85, 12.7, 0.38, C_GOLD)
add_text(sl, "自社開発 OSS 構成のため、ユーザー数無制限・追加ライセンス費ゼロ・社内エンジニアによる改修自在",
         0.45, 6.87, 12.3, 0.35,
         size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)


# ── スライド 5: 機能比較表 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "PowerApps vs opeSchedule 機能比較",
             "主要機能・運用面での詳細比較")

# 表ヘッダー
headers = ["比較項目", "PowerApps (現行)", "opeSchedule (提案)", "優位"]
xs = [0.3, 3.0, 6.8, 10.7]
add_rect(sl, 0.3, 2.05, 12.7, 0.55, C_DARK)
for hdr, x in zip(headers, xs):
    w = xs[xs.index(x)+1] - x if xs.index(x)+1 < len(xs) else 2.5
    add_text(sl, hdr, x + 0.08, 2.1, w - 0.1, 0.45,
             size=13, bold=True, color=C_WHITE)

rows = [
    ("ガントチャート表示",  "基本のみ (Power BI 連携要)",  "ドラッグ & ドロップ, マイルストーン◆",  "提案"),
    ("ユーザー数",         "無制限 (M365 依存)",          "完全無制限・追加費ゼロ",                 "提案"),
    ("ライセンス費",       "約 3,000 円/人/月",           "0 円 (OSS)",                            "提案"),
    ("バージョン管理",     "なし (手動)",                 "自動変更ログ + 手動バージョンUP",         "提案"),
    ("データ保存場所",     "Microsoft クラウド",          "自社管理 (AWS or オンプレ)",              "提案"),
    ("カスタマイズ性",     "Power Fx 限定 / 外注依存",    "Python/React — 社内改修可",              "提案"),
    ("オフライン動作",     "不可",                        "オンプレ構成で社内 LAN 完結可",           "提案"),
    ("既存 M365 統合",     "Teams / SharePoint 連携容易", "REST API 経由で要別途実装",              "現行"),
    ("管理者不要運用",     "Microsoft 管理コンソール",    "サーバー管理者が必要 (最小限)",           "現行"),
    ("初期導入コスト",     "即時利用可",                  "環境構築 1〜2 週間",                     "現行"),
]

row_colors = [C_LIGHT_BG, C_WHITE]
advantage_colors = {"提案": C_ACCENT, "現行": C_POWERAPPS, "同等": C_MUTED}

for i, row in enumerate(rows):
    y = 2.65 + i * 0.44
    add_rect(sl, xs[0], y, xs[-1] - xs[0] + 2.5, 0.44, row_colors[i % 2])
    for j, (cell, x) in enumerate(zip(row, xs)):
        w = xs[j+1] - x if j+1 < len(xs) else 2.5
        if j == 3:
            col = advantage_colors.get(cell, C_MUTED)
            add_rect(sl, x, y, w, 0.44, col)
            add_text(sl, cell, x + 0.05, y + 0.07, w - 0.1, 0.32,
                     size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        else:
            add_text(sl, cell, x + 0.08, y + 0.06, w - 0.12, 0.35,
                     size=11, color=C_DARK)


# ── スライド 6: 運用コスト比較 ────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "運用コスト比較",
             "PowerApps 現行 / AWS 運用 / オンプレ運用 — 年間コスト比較（20名規模想定）")

# 3カラムカード
cards = [
    ("PowerApps (現行)",   C_POWERAPPS, [
        ("ライセンス費",     "3,000 円 × 20 名 × 12 ヶ月"),
        ("小計",             "720,000 円/年"),
        ("追加カスタマイズ", "外注費: 50〜200 万円/件"),
        ("管理コスト",       "M365 管理者工数 (年 5〜10 日)"),
        ("5年累計試算",      "約 360〜700 万円"),
    ]),
    ("AWS 運用 (提案)",    C_AWS, [
        ("EC2 t3.small",     "約 2,200 円/月"),
        ("RDS db.t3.micro",  "約 2,200 円/月"),
        ("ALB + Route53",    "約 2,800 円/月"),
        ("小計",             "約 86,400 円/年"),
        ("5年累計試算",      "約 43〜80 万円 (初期含)"),
    ]),
    ("オンプレ運用 (提案)", C_ONPREM, [
        ("サーバー機器(初期)", "20〜50 万円 (1 回のみ)"),
        ("電気代 + 保守",    "約 10,000 円/月"),
        ("IT人件費(月1h)",   "約 5,000 円/月"),
        ("小計",             "約 180,000 円/年 + 初期費"),
        ("5年累計試算",      "約 90〜130 万円 (機器込)"),
    ]),
]

for i, (title, color, items) in enumerate(cards):
    x = 0.4 + i * 4.3
    add_rect(sl, x, 2.05, 4.05, 0.5, color)
    add_text(sl, title, x + 0.1, 2.1, 3.85, 0.42,
             size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, (label, val) in enumerate(items):
        y = 2.6 + j * 0.72
        bg = C_LIGHT_BG if j % 2 == 0 else C_WHITE
        if label in ("小計", "5年累計試算"):
            bg = RGBColor(0xE8, 0xF4, 0xE8) if "提案" in title else \
                 RGBColor(0xF0, 0xE8, 0xF8) if "現行" in title else \
                 RGBColor(0xE8, 0xF0, 0xF8)
        add_rect(sl, x, y, 4.05, 0.68, bg)
        add_text(sl, label, x + 0.1, y + 0.04, 1.6, 0.3,
                 size=11, bold=(label in ("小計", "5年累計試算")), color=C_MUTED)
        add_text(sl, val, x + 0.1, y + 0.33, 3.8, 0.32,
                 size=11, bold=(label in ("小計", "5年累計試算")), color=C_DARK)

# 矢印
add_rect(sl, 4.55, 4.0, 0.2, 0.25, C_MUTED)
add_text(sl, "vs", 4.46, 3.75, 0.4, 0.35, size=13, bold=True,
         color=C_MUTED, align=PP_ALIGN.CENTER)
add_rect(sl, 8.85, 4.0, 0.2, 0.25, C_MUTED)
add_text(sl, "vs", 8.76, 3.75, 0.4, 0.35, size=13, bold=True,
         color=C_MUTED, align=PP_ALIGN.CENTER)

# 注釈
add_rect(sl, 0.3, 6.88, 12.7, 0.35, C_LIGHT_BG)
add_text(sl,
         "※ AWS 費用は 2026年3月 時点の東京リージョン料金 (1USD=150円換算)。"
         "実際の費用はリソースサイズ・通信量により変動します。"
         "  ※ オンプレは機器償却 5年で試算。",
         0.4, 6.9, 12.4, 0.33,
         size=10, italic=True, color=C_MUTED)


# ── スライド 7: AWS 構成案 ────────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "AWS 構成案",
             "スモールスタート構成 — 可用性・スケーラビリティを確保しつつ低コストで運用")

# 左: 構成図（テキストベース）
add_rect(sl, 0.3, 2.0, 6.0, 5.15, C_LIGHT_BG)
add_rect(sl, 0.3, 2.0, 6.0, 0.42, C_AWS)
add_text(sl, "AWS 構成図（概要）", 0.4, 2.05, 5.8, 0.35,
         size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

infra = [
    ("Route 53",          "カスタムドメイン → ALB への DNS 解決"),
    ("ACM (SSL証明書)",   "HTTPS 化（無料で自動更新）"),
    ("ALB",               "ロードバランサー → EC2 への振り分け"),
    ("EC2 t3.small",      "アプリサーバー (2vCPU / 2GB) — Docker 実行"),
    ("  └ Docker",        "opeSchedule (FastAPI + Vite ビルド済み)"),
    ("RDS db.t3.micro",   "PostgreSQL 15 — Multi-AZ オプション可"),
    ("S3",                "DB バックアップ・ログ保存"),
    ("CloudWatch",        "監視・アラート → メール通知"),
    ("VPC / SG",          "アクセス制御 (443/80 のみ開放)"),
]
for i, (comp, desc) in enumerate(infra):
    y = 2.5 + i * 0.52
    add_rect(sl, 0.35, y, 1.85, 0.47, C_AWS if i in (0,1,2,3,5) else C_DARK)
    add_text(sl, comp, 0.38, y + 0.06, 1.8, 0.37,
             size=10, bold=True, color=C_WHITE)
    add_text(sl, desc, 2.25, y + 0.06, 3.95, 0.4,
             size=10, color=C_DARK)

# 右: スペック・コスト詳細
add_rect(sl, 6.6, 2.0, 6.4, 5.15, C_WHITE)
add_rect(sl, 6.6, 2.0, 6.4, 0.42, C_PRIMARY)
add_text(sl, "月額コスト内訳（概算）", 6.7, 2.05, 6.2, 0.35,
         size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

cost_rows = [
    ("EC2 t3.small",         "~$15",     "~2,200 円"),
    ("RDS db.t3.micro",      "~$15",     "~2,200 円"),
    ("ALB",                  "~$16",     "~2,400 円"),
    ("Route 53",             "~$0.5",    "~75 円"),
    ("S3 (バックアップ)",    "~$1",      "~150 円"),
    ("CloudWatch",           "~$2",      "~300 円"),
    ("データ転送 (出)",      "~$1",      "~150 円"),
    ("合計 (月)",            "~$51",     "~7,600 円"),
    ("合計 (年)",            "~$612",    "~91,800 円"),
]
for i, (item, usd, jpy) in enumerate(cost_rows):
    y = 2.5 + i * 0.52
    is_total = item.startswith("合計")
    bg = C_GOLD if is_total else (C_LIGHT_BG if i % 2 == 0 else C_WHITE)
    add_rect(sl, 6.6, y, 6.4, 0.48, bg)
    tcol = C_DARK if not is_total else C_DARK
    add_text(sl, item, 6.7, y + 0.06, 3.2, 0.38,
             size=11, bold=is_total, color=tcol)
    add_text(sl, usd, 9.9, y + 0.06, 1.2, 0.38,
             size=11, bold=is_total, color=tcol, align=PP_ALIGN.RIGHT)
    add_text(sl, jpy, 11.1, y + 0.06, 1.7, 0.38,
             size=11, bold=is_total, color=tcol, align=PP_ALIGN.RIGHT)

# 特記
add_rect(sl, 0.3, 7.15, 12.7, 0.1, C_PRIMARY)
add_text(sl,
         "スケールアップ時: EC2 → t3.medium / ECS Fargate, RDS → Multi-AZ で対応可。CloudFront 追加で全国高速化も容易。",
         0.3, 7.0, 12.7, 0.28, size=10, italic=True, color=C_MUTED)


# ── スライド 8: オンプレ構成案 ────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "オンプレ構成案",
             "社内 LAN 完結構成 — データを自社管理、インターネット不要で運用可能")

# 左: 構成概要
add_rect(sl, 0.3, 2.0, 6.0, 5.15, C_LIGHT_BG)
add_rect(sl, 0.3, 2.0, 6.0, 0.42, C_ONPREM)
add_text(sl, "オンプレ構成図（概要）", 0.4, 2.05, 5.8, 0.35,
         size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

onprem = [
    ("社内サーバー",        "Ubuntu 22.04 LTS — 既存サーバー転用可"),
    ("  Docker",           "opeSchedule コンテナ実行"),
    ("  PostgreSQL 15",    "Docker コンテナ or ホスト直インストール"),
    ("Nginx (リバプロ)",   "社内 HTTP/HTTPS ルーティング"),
    ("社内 DNS",           "ホスト名解決 (例: gantt.internal)"),
    ("バックアップ",        "cron + pg_dump → NAS/共有ドライブ"),
    ("アクセス制御",        "社内 LAN のみ — ファイアウォールで外部遮断"),
    ("監視",               "Zabbix / Mackerel（任意）— メール通知"),
    ("OS 更新",            "月 1 回 unattended-upgrades 適用推奨"),
]
for i, (comp, desc) in enumerate(onprem):
    y = 2.5 + i * 0.52
    add_rect(sl, 0.35, y, 1.85, 0.47, C_ONPREM if i <= 5 else C_DARK)
    add_text(sl, comp, 0.38, y + 0.06, 1.8, 0.37,
             size=10, bold=True, color=C_WHITE)
    add_text(sl, desc, 2.25, y + 0.06, 3.95, 0.4,
             size=10, color=C_DARK)

# 右: スペック・コスト
add_rect(sl, 6.6, 2.0, 6.4, 5.15, C_WHITE)
add_rect(sl, 6.6, 2.0, 6.4, 0.42, C_ONPREM)
add_text(sl, "コスト内訳", 6.7, 2.05, 6.2, 0.35,
         size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 初期費用
add_rect(sl, 6.6, 2.48, 6.4, 0.35, C_PRIMARY)
add_text(sl, "初期費用（1 回のみ）", 6.7, 2.52, 6.2, 0.3,
         size=12, bold=True, color=C_WHITE)
initial = [
    ("サーバー機器",     "20〜50 万円 (新規購入)"),
    ("  (既存転用)",     "0 円 (社内遊休サーバー利用)"),
    ("OS / ソフトウェア","0 円 (OSS のみ使用)"),
    ("環境構築工数",     "3〜5 日（社内 IT）"),
]
for i, (item, val) in enumerate(initial):
    y = 2.88 + i * 0.44
    bg = C_LIGHT_BG if i % 2 == 0 else C_WHITE
    add_rect(sl, 6.6, y, 6.4, 0.42, bg)
    add_text(sl, item, 6.7, y + 0.06, 3.5, 0.32, size=11, color=C_DARK)
    add_text(sl, val, 10.1, y + 0.06, 2.7, 0.32,
             size=11, color=C_DARK, align=PP_ALIGN.RIGHT)

# ランニング
add_rect(sl, 6.6, 4.65, 6.4, 0.35, C_PRIMARY)
add_text(sl, "ランニングコスト（月額）", 6.7, 4.69, 6.2, 0.3,
         size=12, bold=True, color=C_WHITE)
running = [
    ("電気代（300W 換算）",  "約 2,000 円/月"),
    ("保守 IT 工数（1h/月）", "約 5,000 円/月"),
    ("SSL 証明書",           "0 円 (Let's Encrypt)"),
    ("合計 (月)",            "約 7,000 円/月"),
    ("合計 (年)",            "約 84,000 円/年"),
]
for i, (item, val) in enumerate(running):
    y = 5.05 + i * 0.44
    is_total = item.startswith("合計")
    bg = C_GOLD if is_total else (C_LIGHT_BG if i % 2 == 0 else C_WHITE)
    add_rect(sl, 6.6, y, 6.4, 0.42, bg)
    add_text(sl, item, 6.7, y + 0.06, 3.5, 0.32,
             size=11, bold=is_total, color=C_DARK)
    add_text(sl, val, 10.1, y + 0.06, 2.7, 0.32,
             size=11, bold=is_total, color=C_DARK, align=PP_ALIGN.RIGHT)

add_text(sl,
         "既存サーバー転用の場合、初期コストほぼゼロで導入可能。インターネット接続不要で情報漏洩リスクを最小化。",
         0.3, 7.0, 12.7, 0.28, size=10, italic=True, color=C_MUTED)


# ── スライド 9: コスト試算サマリー ───────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "コスト試算サマリー",
             "5 年間の総保有コスト（TCO）比較 — 20 名規模、現行 PowerApps との差分")

# 棒グラフ風比較
add_rect(sl, 0.3, 2.0, 12.7, 0.42, C_DARK)
add_text(sl, "5 年間 TCO 比較（概算）",
         0.4, 2.05, 12.5, 0.35,
         size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

scenarios = [
    ("PowerApps (現行)  ライセンス費のみ",
     "360 万円",  1.0,  C_POWERAPPS),
    ("PowerApps + カスタマイズ開発 1件",
     "460〜560 万円", 1.27, C_POWERAPPS),
    ("opeSchedule / AWS  (5年)",
     "43〜80 万円",  0.17, C_AWS),
    ("opeSchedule / オンプレ 新規購入",
     "90〜130 万円",  0.25, C_ONPREM),
    ("opeSchedule / オンプレ 既存転用",
     "42〜62 万円",  0.14, C_ACCENT),
]

for i, (label, val, ratio, color) in enumerate(scenarios):
    y = 2.55 + i * 0.84
    bar_w = ratio * 8.5
    add_rect(sl, 0.3, y, bar_w, 0.6, color)
    add_text(sl, label, 0.4, y + 0.1, bar_w - 0.1, 0.42,
             size=11, bold=False, color=C_WHITE)
    add_text(sl, val, 0.3 + bar_w + 0.1, y + 0.1, 3.5, 0.42,
             size=13, bold=True, color=C_DARK)

# 差分ハイライト
add_rect(sl, 0.3, 6.85, 12.7, 0.42, C_ACCENT)
add_text(sl,
         "AWS 構成採用時: PowerApps 比 最大 -83% のコスト削減   "
         "オンプレ(転用)採用時: 最大 -88% のコスト削減",
         0.4, 6.88, 12.4, 0.37,
         size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl,
         "※ PowerApps のカスタマイズ外注費・M365 管理工数は含まず。実際の削減効果はさらに大きい可能性があります。",
         0.3, 7.1, 12.7, 0.22, size=10, italic=True, color=C_MUTED)


# ── スライド 10: 移行・導入スケジュール ──────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "移行・導入スケジュール",
             "段階的アプローチ — 現行 PowerApps と並行運用し、リスクを最小化")

phases = [
    ("Phase 1", "Week 1〜2",  C_PRIMARY,
     ["環境構築（AWS or オンプレ）",
      "Docker デプロイ + DB 初期化",
      "社内アクセス確認 / SSL 設定",
      "テストユーザー 3〜5 名で試験運用"]),
    ("Phase 2", "Week 3〜4",  C_ACCENT,
     ["既存 PowerApps データの CSV エクスポート",
      "opeSchedule へのインポート作業",
      "並行運用開始（両システム稼働）",
      "操作マニュアル配布・簡易研修"]),
    ("Phase 3", "Week 5〜6",  C_GOLD,
     ["全ユーザーへの切り替え展開",
      "PowerApps へのデータ更新停止",
      "1 週間の安定稼働確認",
      "問題なければ PowerApps ライセンス解約手続き"]),
    ("Phase 4", "Week 7〜",   C_DARK,
     ["PowerApps ライセンス費 削減効果 実現",
      "定期バックアップ運用の確立",
      "監視アラート設定（CloudWatch 等）",
      "四半期ごとの機能要望ヒアリング"]),
]

for i, (phase, period, color, tasks) in enumerate(phases):
    x = 0.3 + i * 3.25
    add_rect(sl, x, 2.05, 3.0, 0.52, color)
    add_text(sl, phase, x + 0.05, 2.1, 1.5, 0.44,
             size=16, bold=True, color=C_WHITE)
    add_text(sl, period, x + 1.55, 2.16, 1.35, 0.36,
             size=11, color=C_WHITE, align=PP_ALIGN.RIGHT)
    for j, task in enumerate(tasks):
        y = 2.62 + j * 0.88
        add_rect(sl, x, y, 3.0, 0.82, C_LIGHT_BG if j % 2 == 0 else C_WHITE)
        add_rect(sl, x, y, 0.06, 0.82, color)
        add_text(sl, task, x + 0.14, y + 0.18, 2.75, 0.5,
                 size=11, color=C_DARK, wrap=True)

# 矢印
for ax in [3.35, 6.6, 9.85]:
    add_rect(sl, ax, 2.23, 0.2, 0.2, C_MUTED)
    add_text(sl, "→", ax - 0.05, 2.18, 0.35, 0.35,
             size=18, bold=True, color=C_MUTED)

# ポイント
add_rect(sl, 0.3, 6.88, 12.7, 0.38, C_PRIMARY)
add_text(sl,
         "並行運用期間を設けることでリスクを最小化。"
         "Phase 3 完了後に PowerApps ライセンスを解約し、翌月からコスト削減効果を実現。",
         0.4, 6.91, 12.4, 0.32,
         size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── スライド 11: リスクと対策 ────────────────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "リスクと対策",
             "想定されるリスクと緩和策 — 意思決定前に確認すべき事項")

risks = [
    ("サーバー障害 / 停止",
     "利用者がスケジュールを参照できなくなる",
     C_WARN,
     "AWS: Multi-AZ RDS + ECS で自動復旧\nオンプレ: 定期バックアップ(日次) + 障害時は\n管理者が即日復旧。RPO 1 日 / RTO 4 時間目安"),
    ("セキュリティ侵害",
     "スケジュールデータの漏洩・改ざん",
     C_WARN,
     "HTTPS 強制 / VPN 限定アクセス\nSecurity Group で 443/80 のみ開放\n定期的な OS / ライブラリ更新"),
    ("担当者異動・属人化",
     "保守できる社内エンジニアがいなくなる",
     RGBColor(0xE8, 0x7C, 0x1A),
     "コード全量を GitHub で管理\n設計書・運用手順書を整備済み\nDockerfile で環境再現性を確保"),
    ("データ移行不備",
     "PowerApps から正しくデータが移行されない",
     RGBColor(0xE8, 0x7C, 0x1A),
     "CSV エクスポートで段階移行\n並行運用期間（2 週間）で差異チェック\nロールバック手順を事前に策定"),
    ("ユーザー習熟コスト",
     "新 UI への慣れに時間がかかる",
     C_ACCENT,
     "運用マニュアル (PPTX) 配布済み\n30 分の操作研修で基本習得可能\n問い合わせ窓口を IT 担当者に設定"),
    ("機能不足・要望対応",
     "PowerApps にある機能が一部未実装",
     C_ACCENT,
     "REST API 完備で追加機能開発が容易\n社内エンジニアが対応可能（Python/React）\nフィードバック収集サイクルを四半期で実施"),
]

for i, (risk_title, risk_desc, risk_color, mitigation) in enumerate(risks):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 2.0 + row * 1.68
    add_rect(sl, x, y, 6.2, 0.42, risk_color)
    add_text(sl, risk_title, x + 0.1, y + 0.06, 4.0, 0.32,
             size=13, bold=True, color=C_WHITE)
    add_text(sl, risk_desc, x + 0.1, y + 0.06, 5.8, 0.32,
             size=11, italic=True,
             color=RGBColor(0xFF, 0xFF, 0xCC), align=PP_ALIGN.RIGHT)
    add_rect(sl, x, y + 0.42, 6.2, 1.18, C_LIGHT_BG)
    txb = sl.shapes.add_textbox(
        Inches(x + 0.12), Inches(y + 0.5),
        Inches(5.96), Inches(1.02))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].add_run().text = mitigation
    tf.paragraphs[0].runs[0].font.size = Pt(11)
    tf.paragraphs[0].runs[0].font.color.rgb = C_DARK

add_rect(sl, 0.3, 7.05, 12.7, 0.1, C_MUTED)
add_text(sl, "リスクレベル:  高（赤）/ 中（橙）/ 低（緑）",
         0.3, 7.1, 12.7, 0.22, size=10, italic=True, color=C_MUTED)


# ── スライド 12: 推奨構成と意思決定フロー ────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "推奨構成と意思決定フロー",
             "組織の状況に応じた最適構成の選択ガイド")

# フロー図（テキストボックス群）
add_text(sl, "導入構成 選択フロー", 0.3, 2.0, 12.7, 0.45,
         size=16, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# 質問ノード
questions = [
    (5.5,  2.55, "社内に Linux/Docker を\n運用できるエンジニアがいる？", C_DARK),
    (2.2,  3.95, "既存の遊休サーバーが\n利用可能？", C_DARK),
    (8.8,  3.95, "機密データを\n社外に出したくない？", C_DARK),
]
for x, y, q, col in questions:
    add_rect(sl, x, y, 3.3, 0.8, C_LIGHT_BG)
    add_rect(sl, x, y, 3.3, 0.06, C_PRIMARY)
    add_text(sl, q, x + 0.1, y + 0.1, 3.1, 0.65,
             size=12, color=col, align=PP_ALIGN.CENTER)

# 回答線と結果ノード
add_text(sl, "YES", 2.6, 3.5, 1.0, 0.32,
         size=12, bold=True, color=C_ACCENT)
add_text(sl, "NO", 8.6, 3.5, 0.8, 0.32,
         size=12, bold=True, color=C_WARN)
add_text(sl, "YES", 1.1, 5.0, 1.0, 0.32,
         size=12, bold=True, color=C_ACCENT)
add_text(sl, "NO", 3.3, 5.0, 0.8, 0.32,
         size=12, bold=True, color=C_WARN)
add_text(sl, "YES", 8.0, 5.0, 0.9, 0.32,
         size=12, bold=True, color=C_ACCENT)
add_text(sl, "NO", 10.4, 5.0, 0.8, 0.32,
         size=12, bold=True, color=C_WARN)

# 結果ボックス
results = [
    (0.3,  5.42, "オンプレ\n(既存転用)", C_ACCENT, "初期ほぼ0円\n最安構成"),
    (3.2,  5.42, "オンプレ\n(新規購入)", C_ONPREM, "初期20〜50万\n月額最安"),
    (6.9,  5.42, "オンプレ\n(セキュア)", C_ONPREM, "データ社内完結\n月額最安"),
    (9.8,  5.42, "AWS\n(推奨)", C_AWS, "月額~7,600円\n管理不要"),
    (5.5,  3.95, "AWS\n(代替)", C_AWS, "月額~7,600円\nIT不要"),
]
for x, y, label, color, note in results:
    add_rect(sl, x, y, 2.5, 1.3, color)
    add_text(sl, label, x + 0.05, y + 0.1, 2.4, 0.65,
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, note, x + 0.05, y + 0.75, 2.4, 0.45,
             size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 0.3, 6.85, 12.7, 0.4, C_GOLD)
add_text(sl,
         "本プロジェクトの推奨: AWS 構成（月額 ~7,600 円）— IT 担当者の負荷最小・スケールアップ容易・バックアップ自動化",
         0.4, 6.88, 12.4, 0.35,
         size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)


# ── スライド 13: まとめ・次のアクション ──────────────────────────────────────
sl = prs.slides.add_slide(BLANK)
slide_header(sl, "まとめ・次のアクション",
             "意思決定のポイントと今後の進め方")

# 左: サマリー
add_rect(sl, 0.3, 2.05, 6.0, 0.42, C_PRIMARY)
add_text(sl, "導入メリット サマリー", 0.4, 2.1, 5.8, 0.35,
         size=13, bold=True, color=C_WHITE)

merits = [
    ("コスト削減",
     "PowerApps ライセンス費（年 72 万円）を\n最大 -88% 削減。5 年で 300 万円以上の差。"),
    ("拡張性",
     "ユーザー数増加でコスト増加ゼロ。\n社内エンジニアが機能追加可能。"),
    ("データ主権",
     "スケジュールデータを自社管理。\nM365 解約の影響を受けない。"),
    ("機能優位性",
     "ガントチャート・バージョン管理・\n変更ログ機能が標準搭載。"),
    ("移行リスク小",
     "2 週間の並行運用で安全移行。\nロールバック手順も整備済み。"),
]
for i, (title, desc) in enumerate(merits):
    y = 2.52 + i * 0.92
    add_rect(sl, 0.3, y, 0.55, 0.82, C_ACCENT)
    add_text(sl, str(i+1), 0.3, y + 0.16, 0.55, 0.5,
             size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 0.85, y, 5.45, 0.82, C_LIGHT_BG)
    add_text(sl, title, 0.95, y + 0.04, 1.5, 0.34,
             size=12, bold=True, color=C_PRIMARY)
    add_text(sl, desc, 0.95, y + 0.38, 5.2, 0.4,
             size=11, color=C_DARK)

# 右: 次のアクション
add_rect(sl, 6.6, 2.05, 6.4, 0.42, C_DARK)
add_text(sl, "次のアクション（承認後）", 6.7, 2.1, 6.2, 0.35,
         size=13, bold=True, color=C_WHITE)

actions = [
    ("Week 1",  "構成決定（AWS or オンプレ）\n予算申請・承認"),
    ("Week 2",  "環境構築 + Docker デプロイ\nテストユーザー選定"),
    ("Week 3",  "データ移行テスト + 並行運用開始\nマニュアル配布"),
    ("Week 5",  "全ユーザー切り替え\nPowerApps 更新停止"),
    ("Week 7",  "安定稼働確認\nPowerApps ライセンス解約申請"),
    ("Month 3", "コスト削減効果レポート作成\n継続改善ヒアリング"),
]
for i, (period, action) in enumerate(actions):
    y = 2.52 + i * 0.77
    add_rect(sl, 6.6, y, 1.3, 0.72, C_PRIMARY)
    add_text(sl, period, 6.63, y + 0.14, 1.2, 0.42,
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 7.9, y, 5.1, 0.72,
             C_LIGHT_BG if i % 2 == 0 else C_WHITE)
    add_text(sl, action, 8.0, y + 0.08, 4.9, 0.58,
             size=11, color=C_DARK)

# 決定ボックス
add_rect(sl, 0.3, 7.0, 12.7, 0.45, C_ACCENT)
add_text(sl,
         "ご承認をいただければ、最短 2 週間でパイロット運用を開始できます。"
         "  まずは構成（AWS / オンプレ）の選択をお願いします。",
         0.4, 7.03, 12.4, 0.38,
         size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ── 保存 ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "opeSchedule_proposal.pptx")
prs.save(out)
print(f"[OK] Saved: {out}  ({prs.slides.__len__()} slides)")
